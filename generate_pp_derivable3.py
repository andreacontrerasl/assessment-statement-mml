import os
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
from tkinter import messagebox

def analyze_market_section(df):
    """
    Analyze revenue by market section without converting to USD.
    """
    df['Client'] = df['Client'].astype(str)
    df['Market_Section'] = df['Client'].str[0]
    revenue_by_section = df.groupby('Market_Section')['Transaction'].sum().reset_index()
    revenue_by_section = revenue_by_section.sort_values('Transaction', ascending=False)
    return revenue_by_section

def analyze_geography(df):
    """
    Analyze revenue by country without converting to USD.
    """
    revenue_by_country = df.groupby('Country')['Transaction'].sum().reset_index()
    revenue_by_country = revenue_by_country.sort_values('Transaction', ascending=False)
    return revenue_by_country

def generate_charts(revenue_by_section, revenue_by_country, output_dir):
    sns.set(style="whitegrid")

    # Revenue by Market Section
    plt.figure(figsize=(10, 6))
    sns.barplot(x='Market_Section', y='Transaction', data=revenue_by_section, palette='viridis')
    plt.title('Total Revenue by Market Section')
    plt.xlabel('Market Section')
    plt.ylabel('Total Transactions')
    plt.tight_layout()
    market_section_img_path = os.path.join(output_dir, 'revenue_by_market_section.png')
    plt.savefig(market_section_img_path, dpi=300)
    plt.close()

    # Revenue by Country (Top 20)
    plt.figure(figsize=(12, 8))
    top_countries = revenue_by_country.head(20)
    sns.barplot(x='Transaction', y='Country', data=top_countries, palette='magma')
    plt.title('Total Revenue by Country (Top 20)')
    plt.xlabel('Total Transactions')
    plt.ylabel('Country')
    plt.tight_layout()
    country_img_path = os.path.join(output_dir, 'revenue_by_country.png')
    plt.savefig(country_img_path, dpi=300)
    plt.close()

    return market_section_img_path, country_img_path


def generate_presentation(output_dir, revenue_by_section_img, revenue_by_country_img, revenue_by_section, revenue_by_country):
    """
    Generate a comprehensive PowerPoint presentation with the analysis and graphics.
    """
    # Obtener la sección de mercado y país con mayores ingresos
    top_market_section = revenue_by_section.iloc[0]['Market_Section']
    top_country = revenue_by_country.iloc[0]['Country']
    other_top_countries = ', '.join(revenue_by_country['Country'].iloc[1:4])

    prs = Presentation()

    # Slide 1: Title and Executive Summary
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Financial Analysis Report"
    subtitle.text = "Market Section and Geographical Analysis"
    
    # Executive Summary slide
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Executive Summary"
    content.text = (
        "This report provides an in-depth analysis of revenue by market section and geography. "
        "Key highlights include the identification of the most profitable market sections and countries, "
        "along with strategic recommendations for revenue growth."
    )

    # Slide 2: Market Section Overview
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Market Section Overview"
    content.text = (
        f"The market section '{top_market_section}' stands out as the most lucrative. "
        "A detailed review of the top segments reveals varying revenue patterns that present both opportunities and challenges."
    )

    # Slide 3: Detailed Market Section Analysis
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Market Section Analysis"
    content.text = (
        f"The analysis indicates that market section '{top_market_section}' leads in revenue generation. "
        f"However, sections like {', '.join(revenue_by_section['Market_Section'].iloc[1:4])} also show potential growth opportunities. "
        "Focus on these sections may yield higher returns. Conversely, sections with lower revenues may require "
        "a reevaluation of strategy or resources."
    )

    # Slide 4: Geographical Revenue Overview
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Geographical Revenue Overview"
    content.text = (
        f"The geographical analysis shows that '{top_country}' is the highest revenue-generating country. "
        "This suggests market strength in this region. Countries like "
        f"{other_top_countries} also contribute significantly, indicating a diverse revenue base."
    )

    # Slide 5: Detailed Geographical Analysis
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Geographical Analysis"
    content.text = (
        f"The high revenue in '{top_country}' suggests a strong market presence. "
        "Investment in marketing and operations in this region could further solidify market dominance. "
        "In contrast, regions with lower revenues could represent untapped opportunities or markets that require a different strategic approach."
    )

    # Slide 6: Charts - Market Section and Country
    slide_layout = prs.slide_layouts[5]  # Blank layout for adding images
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Revenue by Market Section"
    slide.shapes.add_picture(revenue_by_section_img, Inches(1), Inches(1.5), width=Inches(8))

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Revenue by Country (Top 20)"
    slide.shapes.add_picture(revenue_by_country_img, Inches(1), Inches(1.5), width=Inches(8))

    # Slide 7: Recommendations and Next Steps
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Recommendations and Next Steps"
    content.text = (
        "1. Prioritize market sections like 'A' and 'B' to maximize profitability.\n"
        "2. Invest in marketing and operational efforts in top-performing countries like Peru.\n"
        "3. Explore potential growth in underperforming sections and regions.\n"
        "4. Conduct further analysis to identify factors driving success in high-revenue segments.\n"
        "5. Develop targeted strategies to enhance revenue in regions with lower market penetration."
    )

    # Slide 8: Conclusions and Acknowledgements
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Conclusions and Acknowledgements"
    content.text = (
        "The analysis provides valuable insights into revenue distribution across market sections and geographies. "
        "Implementing the recommended strategies can enhance revenue growth and optimize market presence. "
        "Thank you for your attention and consideration of these findings."
    )

    # Save the presentation
    presentation_path = os.path.join(output_dir, "Financial_Analysis_Report.pptx")
    prs.save(presentation_path)
    print(f"Presentation saved at: {presentation_path}")

    
def execute_exercise3(final_df):
    # Perform analysis and generate charts without currency conversion
    print("Performing analysis...")
    revenue_by_section = analyze_market_section(final_df)
    revenue_by_country = analyze_geography(final_df)
    
    # Save charts in the Downloads folder
    output_dir = os.path.join(os.path.expanduser('~'), 'Downloads')
    print("Generating charts...")
    revenue_by_section_img, revenue_by_country_img = generate_charts(revenue_by_section, revenue_by_country, output_dir)
    
    # Generate the PowerPoint presentation
    print("Generating PowerPoint presentation...")
    generate_presentation(output_dir, revenue_by_section_img, revenue_by_country_img)
    print(f"Presentation generated in {output_dir}")